import pandas as pd
#import snowflake.connector
import streamlit as st
from streamlit_dynamic_filters import DynamicFilters
from st_aggrid import AgGrid
from st_aggrid.grid_options_builder import GridOptionsBuilder
from st_aggrid import GridUpdateMode, DataReturnMode
import warnings
from yahooquery import Ticker
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
from datetime import date
from PIL import Image
import requests
import os
from io import BytesIO
from langchain.chat_models import ChatOpenAI
from langchain.schema import HumanMessage, SystemMessage
import traceback
import re
import ast

# hide future warnings (caused by st_aggrid)
warnings.simplefilter(action='ignore', category=FutureWarning)

# set page layout and define basic variables
st.set_page_config(layout="wide", page_icon='⚡', page_title="Instant Insight")
path = os.path.dirname(__file__)
today = date.today()


def resize_image(url):
    """function to resize logos while keeping aspect ratio. Accepts URL as an argument and return an image object"""

    # Open the image file
    image = Image.open(requests.get(url, stream=True).raw)

    # if a logo is too high or too wide then make the background container twice as big
    if image.height > 140:
        container_width = 220 * 2
        container_height = 140 * 2

    elif image.width > 220:
        container_width = 220 * 2
        container_height = 140 * 2
    else:
        container_width = 220
        container_height = 140

    # Create a new image with the same aspect ratio as the original image
    new_image = Image.new('RGBA', (container_width, container_height))

    # Calculate the position to paste the image so that it is centered
    x = (container_width - image.width) // 2
    y = (container_height - image.height) // 2

    # Paste the image onto the new image
    new_image.paste(image, (x, y))
    return new_image


def add_image(slide, image, left, top, width):
    """function to add an image to the PowerPoint slide and specify its position and width"""
    slide.shapes.add_picture(image, left=left, top=top, width=width)


def replace_text(replacements, slide):
    """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text


def get_stock(ticker, period, interval):
    """function to get stock data from Yahoo Finance. Takes ticker, period and interval as arguments and returns a DataFrame"""
    hist = ticker.history(period=period, interval=interval)
    hist = hist.reset_index()
    # capitalize column names
    hist.columns = [x.capitalize() for x in hist.columns]
    return hist


def plot_graph(df, x, y, title, name):
    """function to plot a line graph. Takes DataFrame, x and y axis, title and name as arguments and returns a Plotly figure"""
    fig = px.line(df, x=x, y=y, template='simple_white',
                  title='<b>{} {}</b>'.format(name, title))
    fig.update_traces(line_color='#A27D4F')
    fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
    return fig


def peers_plot(df, name, metric):
    """function to plot a bar chart with peers. Takes DataFrame, name, metric and ticker as arguments and returns a Plotly figure"""

    # drop rows with missing metrics
    df.dropna(subset=[metric], inplace=True)

    df_sorted = df.sort_values(metric, ascending=False)

    # iterate over the labels and add the colors to the color mapping dictionary, hightlight the selected ticker
    color_map = {}
    for label in df_sorted['Company Name']:
        if label == name:
            color_map[label] = '#A27D4F'
        else:
            color_map[label] = '#D9D9D9'

    fig = px.bar(df_sorted, y='Company Name', x=metric, template='simple_white', color='Company Name',
                 color_discrete_map=color_map,
                 orientation='h',
                 title='<b>{} {} vs Peers FY22</b>'.format(name, metric))
    fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)', showlegend=False, yaxis_title='')
    return fig


def esg_plot(name, df):
    # Define colors for types
    colors = {name: '#A27D4F', 'Peer Group': '#D9D9D9'}

    # Creating the bar chart
    fig = go.Figure()
    for type in df['Type'].unique():
        fig.add_trace(go.Bar(
            x=df[df['Type'] == type]['variable'],
            y=df[df['Type'] == type]['value'],
            name=type,
            text=df[df['Type'] == type]['value'],
            textposition='outside',
            marker_color=colors[type]
        ))
    fig.update_layout(
        height=700,
        width=1000,
        barmode='group',
        title="ESG Score vs Peers Average",
        xaxis_title="",
        yaxis_title="Score",
        legend_title="Type",
        xaxis=dict(tickangle=0),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)')
    return fig


def get_financials(df, col_name, metric_name):
    """function to get financial metrics from a DataFrame. Takes DataFrame, column name and metric name as arguments and returns a DataFrame"""
    metric = df.loc[:, ['asOfDate', col_name]]
    metric_df = pd.DataFrame(metric).reset_index()
    metric_df.columns = ['Symbol', 'Year', metric_name]

    return metric_df


def generate_gpt_response(gpt_input, max_tokens, api_key, llm_model):
    """function to generate a response from GPT-3. Takes input and max tokens as arguments and returns a response"""
    # Create an instance of the OpenAI class
    chat = ChatOpenAI(openai_api_key=api_key, model=llm_model,
                      temperature=0, max_tokens=max_tokens)

    # Generate a response from the model
    response = chat.predict_messages(
        [SystemMessage(content='You are a helpful expert in finance, market and company research.'
                               'You also have exceptional skills in selling B2B software products.'),
         HumanMessage(
             content=gpt_input)])

    return response.content.strip()


def dict_from_string(response):
    """function to parse GPT response with competitors tickers and convert it to a dict"""
    # Find a substring that starts with '{' and ends with '}', across multiple lines
    match = re.search(r'\{.*?\}', response, re.DOTALL)

    dictionary = None
    if match:
        try:
            # Try to convert substring to dict
            dictionary = ast.literal_eval(match.group())
        except (ValueError, SyntaxError):
            # Not a dictionary
            return None
    return dictionary


def extract_comp_financials(tkr, comp_name, dict):
    """function to extract financial metrics for competitors. Takes a ticker as an argument and appends financial metrics to dict"""
    ticker = Ticker(tkr)
    income_df = ticker.income_statement(frequency='a', trailing=False)

    subset = income_df.loc[:, ['asOfDate', 'TotalRevenue', 'SellingGeneralAndAdministration']].reset_index()

    # keep only 2022 data
    subset = subset[subset['asOfDate'].dt.year == 2022].sort_values(by='asOfDate', ascending=False).head(1)

    # get values
    total_revenue = subset['TotalRevenue'].values[0]
    sg_and_a = subset['SellingGeneralAndAdministration'].values[0]

    # calculate sg&a as a percentage of total revenue
    sg_and_a_pct = round(sg_and_a / total_revenue * 100, 2)

    # add values to dictionary
    dict[comp_name]['Total Revenue'] = total_revenue
    dict[comp_name]['SG&A % Of Revenue'] = sg_and_a_pct


def convert_to_nested_dict(input_dict, nested_key):
    """function to convert a dictionary to a nested dictionary. Takes a dictionary and a nested key as arguments and returns a dictionary"""
    output_dict = {}
    for key, value in input_dict.items():
        output_dict[key] = {nested_key: value}
    return output_dict


def shorten_summary(text):
    # Split the text into sentences using a regular expression pattern
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)

    # Return the first two sentences or less if there are fewer sentences
    sen = ' '.join(sentences[:2])

    # if the summary is less than 350 characters, return the summary
    if len(sen) <= 400:
        return sen
    else:
        # if the summary is more than 350 characters, return the first 350 characters and truncate the last word
        truncated_sen = text[:400].rsplit(' ', 1)[0] + '...'
        return truncated_sen


def peers_summary(df, selected_ticker):
    df = df[df['Ticker'] != selected_ticker]

    for tkr in df['Ticker']:
        try:
            profile = Ticker(tkr).asset_profile
            summary = profile[tkr]['longBusinessSummary']
            website = profile[tkr]['website']

            # keep only the first two sentences of the summary
            short_summary = shorten_summary(summary)
            logo_url = 'https://logo.clearbit.com/' + website

            # append short summary and logo_url to the df
            df.loc[df['Ticker'] == tkr, 'Summary'] = short_summary
            df.loc[df['Ticker'] == tkr, 'Logo'] = logo_url
        except:
            continue
    # drop rows with missing summary
    df = df.dropna(subset=['Summary'])
    return df


def fix_text_capitalization(text):
    fixed_text = text.lower().capitalize()
    return fixed_text


def replace_multiple_symbols(string):
    """function to fix description from yahoo finance, sometimes it has multiple dots at the end of the string"""
    string = string.replace(':', '')
    pattern = r'\.{2,}$'  # Matches two or more consecutive dots at the end of the string
    replacement = '.'

    # Check if the string ends with multiple symbols
    if re.search(pattern, string):
        # Replace multiple symbols with a single symbol
        string = re.sub(pattern, replacement, string)
    return string


def no_data_plot():
    """plot to return when there is no data available"""
    # Create a blank figure with a transparent background
    fig = go.Figure()

    # Add a text annotation for "NO DATA AVAILABLE" at the center of the plot
    fig.add_annotation(
        x=0.5,
        y=0.5,
        xref='paper',
        yref='paper',
        text='NO DATA AVAILABLE',
        showarrow=False,
        font=dict(size=26, color='black'),
    )

    # Customize layout to have a transparent background
    fig.update_layout(
        paper_bgcolor='rgba(0,0,0,0)',  # Transparent background
        plot_bgcolor='rgba(0,0,0,0)',  # Transparent plot area
        xaxis_showgrid=False,  # Hide x-axis gridlines
        yaxis_showgrid=False,  # Hide y-axis gridlines
        xaxis=dict(visible=False),  # Hide x-axis labels and ticks
        yaxis=dict(visible=False),  # Hide y-axis labels and ticks
    )
    return fig


#conn = st.connection("snowflake")
#df = conn.query("SELECT * from prospects LIMIT 1000;", ttl=600)
df = pd.read_csv('prospects.csv')

# Fix column names. Replace underscore with space, lowercase column names, and capitalize first words
df.columns = df.columns.str.replace('_', ' ').str.lower().str.title()

with st.sidebar:
    openai_key = st.text_input(label="Your OpenAI API key", help="Your API key is not stored anywhere")
    llm_model = st.selectbox(label="Choose a model", options=["gpt-3.5-turbo", "gpt-4-turbo", "gpt-4", "gpt-4o"])

# create sidebar filters
st.sidebar.write('**Use filters to select prospects** 👇')

# display dynamic multi select filters
dynamic_filters = DynamicFilters(df, filters=['Sector', 'Industry', 'Prospect Status', 'Product'])
dynamic_filters.display_filters(location='sidebar')
df_filtered = dynamic_filters.filter_df()
